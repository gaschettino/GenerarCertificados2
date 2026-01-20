import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import tempfile
import subprocess
import shutil
import re

# =========================
# Utilidades
# =========================
def css_font(f):
    """Evita problemas de fuentes con espacios en HTML"""
    return f.replace("'", "").replace('"', "")

# =========================
# Configuración de la página
# =========================
st.set_page_config(page_title="Generador de Certificados", layout="centered")

st.title("Generador de Certificados")
st.write("Cualquier consulta enviar mail a gaschettino@garrahan.gov.ar")

# =========================
# Subida de archivos (2 columnas)
# =========================
st.subheader("Archivos")

col1, col2 = st.columns(2)

with col1:
    uploaded_template = st.file_uploader(
        "Template del certificado (.pptx)", type=["pptx"]
    )

with col2:
    uploaded_excel = st.file_uploader(
        "Listado de asistentes (.xlsx)", type=["xlsx"]
    )

# =========================
# DNI
# =========================
st.subheader("Contenido del certificado")

incluye_dni = st.checkbox("El certificado incluye DNI")

if incluye_dni:
    st.info(
        "El Excel debe tener una columna 'DNI' "
        "y el template debe incluir el texto 'Numero de DNI'."
    )

st.divider()

# =========================
# Fuentes disponibles
# =========================
fuentes_disponibles = [
    "DejaVu Sans",
    "DejaVu Serif",
    "Liberation Sans",
    "Liberation Serif",
    "Arial",
    "Times New Roman",
    "Calibri"
]

colores_predef = {
    "Negro": (0, 0, 0),
    "Azul": (0, 0, 180),
    "Rojo": (180, 0, 0),
    "Verde": (0, 140, 0),
    "Gris": (90, 90, 90)
}

# =========================
# Formato Nombre
# =========================
st.subheader("Formato del nombre y apellido")

col_nom_1, col_nom_2 = st.columns(2)

with col_nom_1:
    fuente_nombre = st.selectbox(
        "Fuente (Nombre)",
        fuentes_disponibles,
        index=0
    )

    tamaño_nombre = st.number_input(
        "Tamaño de letra (Nombre)",
        min_value=8,
        max_value=60,
        value=25,
        step=1
    )

with col_nom_2:
    st.markdown("Color del nombre  \n👉 [Ver códigos](https://htmlcolorcodes.com/)")

    modo_color_nombre = st.radio(
        "Modo de color",
        ["predefinido", "rgb", "hex"],
        horizontal=True,
        key="color_nombre"
    )

    r_nom, g_nom, b_nom = 0, 0, 0

    if modo_color_nombre == "predefinido":
        c = st.selectbox("Color", colores_predef.keys(), key="c_nom")
        r_nom, g_nom, b_nom = colores_predef[c]

    elif modo_color_nombre == "rgb":
        rgb = st.text_input("RGB (ej: 0,0,0)", key="rgb_nom")
        try:
            r_nom, g_nom, b_nom = map(int, rgb.split(","))
        except:
            pass

    elif modo_color_nombre == "hex":
        hx = st.text_input("HEX (ej: #000000)", key="hex_nom")
        if re.match(r"^#([A-Fa-f0-9]{6})$", hx):
            r_nom = int(hx[1:3], 16)
            g_nom = int(hx[3:5], 16)
            b_nom = int(hx[5:7], 16)

# =========================
# Formato DNI
# =========================
if incluye_dni:
    st.divider()
    st.subheader("Formato del DNI")

    col_dni_1, col_dni_2 = st.columns(2)

    with col_dni_1:
        fuente_dni = st.selectbox(
            "Fuente (DNI)",
            fuentes_disponibles,
            index=0,
            key="fuente_dni"
        )

        tamaño_dni = st.number_input(
            "Tamaño de letra (DNI)",
            min_value=8,
            max_value=40,
            value=14,
            step=1
        )

    with col_dni_2:
        st.markdown("Color del DNI")

        modo_color_dni = st.radio(
            "Modo de color",
            ["predefinido", "rgb", "hex"],
            horizontal=True,
            key="color_dni"
        )

        r_dni, g_dni, b_dni = 0, 0, 0

        if modo_color_dni == "predefinido":
            c = st.selectbox("Color", colores_predef.keys(), key="c_dni")
            r_dni, g_dni, b_dni = colores_predef[c]

        elif modo_color_dni == "rgb":
            rgb = st.text_input("RGB (ej: 0,0,0)", key="rgb_dni")
            try:
                r_dni, g_dni, b_dni = map(int, rgb.split(","))
            except:
                pass

        elif modo_color_dni == "hex":
            hx = st.text_input("HEX (ej: #000000)", key="hex_dni")
            if re.match(r"^#([A-Fa-f0-9]{6})$", hx):
                r_dni = int(hx[1:3], 16)
                g_dni = int(hx[3:5], 16)
                b_dni = int(hx[5:7], 16)

st.divider()

# =========================
# Vista previa
# =========================
st.subheader("Vista previa")

preview_html = f"""<div style="display:flex; gap:40px; align-items:flex-start; margin-top:10px;">
<div style="font-family:'{css_font(fuente_nombre)}', sans-serif;
font-size:{tamaño_nombre}px;
font-weight:bold;
font-style:italic;
color:rgb({r_nom},{g_nom},{b_nom});">
Nombre y Apellido
</div>"""

if incluye_dni:
    preview_html += f"""<div style="font-family:'{css_font(fuente_dni)}', sans-serif;
font-size:{tamaño_dni}px;
color:rgb({r_dni},{g_dni},{b_dni});">
DNI 12.345.678
</div>"""

preview_html += "</div>"

st.markdown(preview_html, unsafe_allow_html=True)

st.caption(
    "La vista previa es orientativa. "
    "El PDF final usa las fuentes disponibles en el servidor."
)


st.divider()

# =========================
# Conversión PPTX → PDF
# =========================
def convert_to_pdf(input_pptx, output_dir):
    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, input_pptx],
        check=True
    )
    os.remove(input_pptx)

# =========================
# Procesamiento
# =========================
if uploaded_template and uploaded_excel:
    if st.button("🚀 Generar certificados"):
        with st.spinner("Generando certificados..."):
            with tempfile.TemporaryDirectory() as tmpdir:

                template_path = os.path.join(tmpdir, "template.pptx")
                with open(template_path, "wb") as f:
                    f.write(uploaded_template.read())

                df = pd.read_excel(uploaded_excel)
                df.columns = df.columns.str.strip().str.title()

                if "Nombre" not in df.columns or "Apellido" not in df.columns:
                    st.error("El Excel debe tener columnas Nombre y Apellido.")
                    st.stop()

                if incluye_dni and "Dni" not in df.columns:
                    st.error("Marcaste DNI pero falta la columna DNI.")
                    st.stop()

                df["Nombre y apellido"] = (
                    df["Apellido"].astype(str) + " " + df["Nombre"].astype(str)
                ).str.title()

                out = os.path.join(tmpdir, "Certificados")
                os.makedirs(out, exist_ok=True)

                for _, row in df.iterrows():
                    prs = Presentation(template_path)

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for p in shape.text_frame.paragraphs:
                                    for run in p.runs:
                                        if "Nombre y apellido" in run.text:
                                            run.text = row["Nombre y apellido"]
                                            run.font.name = fuente_nombre
                                            run.font.size = Pt(tamaño_nombre)
                                            run.font.bold = True
                                            run.font.italic = True
                                            run.font.color.rgb = RGBColor(r_nom, g_nom, b_nom)

                                        if incluye_dni and "Numero de DNI" in run.text:
                                            run.text = str(row["Dni"])
                                            run.font.name = fuente_dni
                                            run.font.size = Pt(tamaño_dni)
                                            run.font.color.rgb = RGBColor(r_dni, g_dni, b_dni)

                                    p.alignment = PP_ALIGN.CENTER

                    fname = row["Nombre y apellido"].replace(" ", "_")
                    pptx = os.path.join(out, f"Certificado_{fname}.pptx")
                    prs.save(pptx)
                    convert_to_pdf(pptx, out)

                zip_path = os.path.join(tmpdir, "certificados.zip")
                shutil.make_archive(zip_path.replace(".zip", ""), "zip", out)

                with open(zip_path, "rb") as f:
                    st.download_button(
                        "Descargar certificados",
                        f,
                        "certificados.zip",
                        "application/zip"
                    )
